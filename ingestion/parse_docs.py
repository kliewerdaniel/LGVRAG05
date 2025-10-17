"""
Document parsing module for various file formats.

This module provides functionality to parse PDF, Markdown, Word documents,
Excel files, PowerPoint presentations, and plain text files.
"""

import os
import logging
import sys
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import re

# Add parent directory to path for imports
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# Document processing libraries - Updated for Python 3.13 compatibility
try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import markdown
    MARKDOWN_AVAILABLE = True
except ImportError:
    MARKDOWN_AVAILABLE = False

from config import get_config

logger = logging.getLogger(__name__)


class DocumentParser:
    """
    A comprehensive document parser that supports multiple file formats.

    This class handles parsing of various document types and provides
    text chunking functionality for further processing.
    """

    def __init__(self, config_path: Optional[str] = None):
        """
        Initialize the document parser.

        Args:
            config_path: Path to configuration file
        """
        self.config = get_config(config_path)
        self.supported_extensions = set(self.config.documents.supported_formats)

    def parse_file(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Parse a single document file.

        Args:
            file_path: Path to the document file

        Returns:
            List of document chunks with metadata

        Raises:
            ValueError: If file format is not supported
            FileNotFoundError: If file doesn't exist
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        file_extension = Path(file_path).suffix.lower()

        if file_extension not in self.supported_extensions:
            raise ValueError(f"Unsupported file format: {file_extension}")

        logger.info(f"Parsing document: {file_path}")

        try:
            if file_extension == '.pdf':
                return self._parse_pdf(file_path)
            elif file_extension == '.md':
                return self._parse_markdown(file_path)
            elif file_extension == '.txt':
                return self._parse_text(file_path)
            elif file_extension == '.docx':
                return self._parse_docx(file_path)
            elif file_extension == '.xlsx':
                return self._parse_xlsx(file_path)
            elif file_extension == '.pptx':
                return self._parse_pptx(file_path)
            else:
                raise ValueError(f"Unsupported format: {file_extension}")

        except Exception as e:
            logger.error(f"Error parsing {file_path}: {str(e)}")
            raise

    def parse_directory(self, directory_path: str) -> List[Dict[str, Any]]:
        """
        Parse all supported documents in a directory.

        Args:
            directory_path: Path to directory containing documents

        Returns:
            List of all document chunks from all files
        """
        if not os.path.exists(directory_path):
            raise FileNotFoundError(f"Directory not found: {directory_path}")

        all_chunks = []

        for file_path in Path(directory_path).rglob('*'):
            if file_path.suffix.lower() in self.supported_extensions:
                try:
                    chunks = self.parse_file(str(file_path))
                    all_chunks.extend(chunks)
                    logger.info(f"Parsed {file_path}: {len(chunks)} chunks")
                except Exception as e:
                    logger.warning(f"Failed to parse {file_path}: {str(e)}")
                    continue

        logger.info(f"Total chunks parsed: {len(all_chunks)}")
        return all_chunks

    def _parse_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """Parse PDF file using PyPDF2."""
        try:
            if PYPDF2_AVAILABLE:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)

                    # Extract text from all pages
                    text_content = ""
                    for page in pdf_reader.pages:
                        text_content += page.extract_text() + "\n"

                    return self._create_chunks(text_content, file_path, "pdf")
            else:
                raise ImportError("PyPDF2 not available for PDF parsing")

        except Exception as e:
            logger.error(f"Failed to parse PDF: {str(e)}")
            # Fallback: try to extract as plain text
            try:
                with open(file_path, 'rb') as f:
                    content = f.read()
                    text_content = content.decode('utf-8', errors='ignore')
                    return self._create_chunks(text_content, file_path, "pdf")
            except Exception as e2:
                logger.error(f"Failed to extract text from PDF: {str(e2)}")
                raise

    def _parse_markdown(self, file_path: str) -> List[Dict[str, Any]]:
        """Parse Markdown file."""
        try:
            if MARKDOWN_AVAILABLE:
                with open(file_path, 'r', encoding='utf-8') as f:
                    md_content = f.read()

                # Convert markdown to HTML first, then extract text
                html_content = markdown.markdown(md_content)
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(html_content, 'html.parser')
                text_content = soup.get_text()
            else:
                # Fallback: just read as plain text
                with open(file_path, 'r', encoding='utf-8') as f:
                    text_content = f.read()

            return self._create_chunks(text_content, file_path, "markdown")

        except Exception as e:
            logger.error(f"Failed to parse markdown: {str(e)}")
            raise

    def _parse_text(self, file_path: str) -> List[Dict[str, Any]]:
        """Parse plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text_content = f.read()

            return self._create_chunks(text_content, file_path, "text")

        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text_content = f.read()
                    return self._create_chunks(text_content, file_path, "text")
                except UnicodeDecodeError:
                    continue

            raise ValueError(f"Could not decode text file: {file_path}")

        except Exception as e:
            logger.error(f"Failed to parse text file: {str(e)}")
            raise

    def _parse_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """Parse Word document."""
        try:
            if DOCX_AVAILABLE:
                doc = Document(file_path)
                text_content = ""

                # Extract text from all paragraphs
                for paragraph in doc.paragraphs:
                    text_content += paragraph.text + "\n"

                # Extract text from tables
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            text_content += cell.text + " "
                        text_content += "\n"

                return self._create_chunks(text_content, file_path, "docx")
            else:
                raise ImportError("python-docx not available for DOCX parsing")

        except Exception as e:
            logger.error(f"Failed to parse DOCX: {str(e)}")
            raise

    def _parse_xlsx(self, file_path: str) -> List[Dict[str, Any]]:
        """Parse Excel file."""
        try:
            if OPENPYXL_AVAILABLE:
                workbook = load_workbook(filename=file_path, read_only=True, data_only=True)
                text_content = ""

                # Extract text from all sheets
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text_content += f"\n--- Sheet: {sheet_name} ---\n"

                    # Read all rows and cells
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                        if row_text.strip():
                            text_content += row_text + "\n"

                return self._create_chunks(text_content, file_path, "xlsx")
            else:
                raise ImportError("openpyxl not available for XLSX parsing")

        except Exception as e:
            logger.error(f"Failed to parse XLSX: {str(e)}")
            raise

    def _parse_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """Parse PowerPoint presentation."""
        try:
            if PPTX_AVAILABLE:
                presentation = Presentation(file_path)
                text_content = ""

                # Extract text from all slides
                for slide_number, slide in enumerate(presentation.slides, 1):
                    text_content += f"\n--- Slide {slide_number} ---\n"

                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content += shape.text + "\n"

                        # Extract text from tables
                        if hasattr(shape, "table"):
                            for row in shape.table.rows:
                                row_text = " | ".join([cell.text for cell in row.cells])
                                text_content += row_text + "\n"

                return self._create_chunks(text_content, file_path, "pptx")
            else:
                raise ImportError("python-pptx not available for PPTX parsing")

        except Exception as e:
            logger.error(f"Failed to parse PPTX: {str(e)}")
            raise

    def _create_chunks(self, text: str, file_path: str, file_type: str) -> List[Dict[str, Any]]:
        """
        Create text chunks from document content.

        Args:
            text: Full text content of the document
            file_path: Path to the source file
            file_type: Type of the file (pdf, markdown, etc.)

        Returns:
            List of text chunks with metadata
        """
        if not text.strip():
            return []

        chunk_size = self.config.documents.chunk_size
        chunk_overlap = self.config.documents.chunk_overlap

        # Clean and normalize text
        text = self._clean_text(text)

        # Split into words
        words = text.split()
        chunks = []

        for i in range(0, len(words), chunk_size - chunk_overlap):
            chunk_words = words[i:i + chunk_size]
            chunk_text = " ".join(chunk_words)

            if len(chunk_text.strip()) < 50:  # Skip very small chunks
                continue

            chunk_metadata = {
                "text": chunk_text,
                "source": file_path,
                "file_type": file_type,
                "chunk_id": f"{Path(file_path).stem}_chunk_{len(chunks)}",
                "start_word": i,
                "end_word": min(i + chunk_size, len(words)),
                "word_count": len(chunk_words),
                "char_count": len(chunk_text)
            }

            chunks.append(chunk_metadata)

        return chunks

    def _clean_text(self, text: str) -> str:
        """
        Clean and normalize text content.

        Args:
            text: Raw text content

        Returns:
            Cleaned text
        """
        # Remove excessive whitespace
        text = re.sub(r'\n\s*\n\s*\n', '\n\n', text)
        text = re.sub(r'[ \t]+', ' ', text)

        # Remove control characters except newlines and tabs
        text = re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f\x7f]', '', text)

        # Normalize quotes
        text = re.sub(r'["\u201c\u201d]', '"', text)  # Normalize double quotes
        text = re.sub(r'[\'\u2018\u2019]', "'", text)  # Normalize single quotes

        return text.strip()


def main():
    """Example usage of the DocumentParser."""
    logging.basicConfig(level=logging.INFO)

    parser = DocumentParser()

    # Parse a single file
    try:
        chunks = parser.parse_file("documents/sample.md")
        print(f"Parsed {len(chunks)} chunks from sample.md")

        if chunks:
            print("First chunk preview:")
            print(f"Text: {chunks[0]['text'][:200]}...")
            print(f"Metadata: {chunks[0]['chunk_id']}")

    except Exception as e:
        print(f"Error: {e}")


if __name__ == "__main__":
    main()
